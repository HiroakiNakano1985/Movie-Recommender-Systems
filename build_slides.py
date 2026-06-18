# -*- coding: utf-8 -*-
"""Build the individual-project slide deck (technical challenges, method
comparison, final remarks) for the movie recommender prototype.

Run (after `python main.py` has produced results/figures):
    python build_slides.py
Output: Movie_Recommender_Individual_Project.pptx
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

ROOT = Path(__file__).resolve().parent
FIG = ROOT / "results" / "figures"
OUT = ROOT / "Movie_Recommender_Individual_Project.pptx"

# ---- palette (matches the course evaluation deck) ----
DARK = RGBColor(0x1F, 0x2A, 0x44)
ACCENT = RGBColor(0x2E, 0x6F, 0xF2)
LIGHT = RGBColor(0xF2, 0xF5, 0xFA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREY = RGBColor(0x55, 0x5E, 0x6E)
WIN = RGBColor(0x1B, 0x7F, 0x4B)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def add_slide():
    return prs.slides.add_slide(BLANK)


def add_box(slide, l, t, w, h):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    return tf


def set_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def title_bar(slide, text):
    bar = slide.shapes.add_shape(1, 0, 0, SW, Inches(1.05))
    bar.fill.solid(); bar.fill.fore_color.rgb = DARK
    bar.line.fill.background()
    tf = bar.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.45); tf.margin_top = Inches(0.18)
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(28); p.font.bold = True; p.font.color.rgb = WHITE


def bullets(tf, items, size=18, space=8):
    first = True
    for txt, lvl, *rest in items:
        bold = rest[0] if rest else False
        color = rest[1] if len(rest) > 1 else DARK
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = txt
        p.level = lvl
        p.font.size = Pt(size - lvl * 2)
        p.font.bold = bold
        p.font.color.rgb = color
        p.space_after = Pt(space)


def style_table(table, header_fill=ACCENT, font=13):
    for r, row in enumerate(table.rows):
        for c, cell in enumerate(row.cells):
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER
                for run in p.runs:
                    run.font.size = Pt(font)
                    if r == 0:
                        run.font.bold = True
                        run.font.color.rgb = WHITE
                    else:
                        run.font.color.rgb = DARK
            if r == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = header_fill
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if r % 2 else LIGHT


def make_table(slide, l, t, w, h, data, header_fill=ACCENT, font=13):
    rows, cols = len(data), len(data[0])
    gt = slide.shapes.add_table(rows, cols, l, t, w, h).table
    for r in range(rows):
        for c in range(cols):
            gt.cell(r, c).text = str(data[r][c])
    style_table(gt, header_fill, font)
    return gt


def add_image(slide, name, l, t, w=None, h=None):
    path = FIG / name
    if path.exists():
        slide.shapes.add_picture(str(path), l, t, width=w, height=h)


def caption(slide, text, l, t, w, color=GREY, size=13):
    tf = add_box(slide, l, t, w, Inches(0.5))
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(size); p.font.italic = True; p.font.color.rgb = color


# ============================================================ 1. TITLE
s = add_slide(); set_bg(s, DARK)
tf = add_box(s, Inches(0.9), Inches(2.1), Inches(11.5), Inches(2.4))
p = tf.paragraphs[0]; p.text = "Movie Recommender Systems"
p.font.size = Pt(46); p.font.bold = True; p.font.color.rgb = WHITE
p = tf.add_paragraph(); p.text = "Building & comparing a recommender prototype — Individual Project"
p.font.size = Pt(24); p.font.color.rgb = RGBColor(0x9C, 0xB6, 0xE8)
tf = add_box(s, Inches(0.95), Inches(5.0), Inches(11.5), Inches(1.6))
for i, (txt, sz) in enumerate([
    ("Recommender Systems — Prof. Marc Torrens, Esade", 18),
    ("Hiroaki Nakano", 18),
    ("June 2026", 14),
]):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = txt; p.font.size = Pt(sz)
    p.font.color.rgb = RGBColor(0xC7, 0xD2, 0xEC)

# ============================================================ 2. GOAL
s = add_slide(); set_bg(s, WHITE); title_bar(s, "Project Goal & Approach")
tf = add_box(s, Inches(0.6), Inches(1.3), Inches(12.1), Inches(5.6))
bullets(tf, [
    ("One prototype, built up method by method on a single dataset.", 0, True, ACCENT),
    ("Movie track — MovieLens Latest Small (100k ratings, 610 users, 9.7k movies).", 1),
    ("Implement and compare seven recommenders on identical data & splits.", 0, True, ACCENT),
    ("Non-personalized baselines, content-based, item-item & user-user CF, matrix factorization.", 1),
    ("Think beyond accuracy.", 0, True, ACCENT),
    ("Also measure coverage, novelty, diversity and popularity bias — accuracy alone is misleading.", 1),
    ("Ship a usable prototype.", 0, True, ACCENT),
    ("Interactive Streamlit app: explore a user, compare algorithms, read 'why' explanations.", 1),
], size=20, space=11)

# ============================================================ 3. DATA + EDA
s = add_slide(); set_bg(s, WHITE); title_bar(s, "Dataset & EDA")
tf = add_box(s, Inches(0.5), Inches(1.2), Inches(6.3), Inches(5.8))
bullets(tf, [
    ("MovieLens Latest Small (GroupLens)", 0, True, ACCENT),
    ("610 users · 9,724 rated movies · 100,836 ratings", 1),
    ("Ratings 0.5–5.0, mean 3.5, mode 4.0 (positivity bias)", 1),
    ("Matrix is 98.3% sparse", 0, True, ACCENT),
    ("Normal — even dense — by recommender standards", 1),
    ("Sparsity is asymmetric", 0, True, ACCENT),
    ("Users are thick: median 70 ratings each (min 20)", 1),
    ("Items are thin: median 3 ratings; heavy long tail", 1),
    ("~20% of movies hold 80% of all ratings", 1, True, GREY),
], size=18, space=8)
add_image(s, "long_tail.png", Inches(7.0), Inches(1.5), w=Inches(6.0))
caption(s, "Long tail: a small head of movies dominates ratings.", Inches(7.0), Inches(5.9), Inches(6.0))

# ============================================================ 4. SPARSITY STRATEGY
s = add_slide(); set_bg(s, WHITE); title_bar(s, "Design Challenge: Sparsity")
tf = add_box(s, Inches(0.6), Inches(1.3), Inches(12.1), Inches(5.6))
bullets(tf, [
    ("Each method handles the thin item axis differently:", 0, True),
    ("Matrix factorization — learns dense latent vectors from observed ratings; built for sparsity.", 0, True, ACCENT),
    ("Content-based — uses item metadata, not co-ratings; works even for long-tail items.", 0, True, ACCENT),
    ("Neighbourhood CF — needs safeguards, or tail items produce noisy similarities:", 0, True, ACCENT),
    ("min_support: drop items rated by < 5 users from the neighbourhood model.", 1),
    ("shrinkage: sim_adj = n_ij / (n_ij + λ) · sim — distrust similarities from few co-ratings.", 1),
    ("adjusted cosine: mean-center each user's ratings before similarity.", 1),
    ("Lesson: standard methods DO work at 98% sparsity — the trick is matching method to the data's weak axis.", 0, True, WIN),
], size=18, space=9)

# ============================================================ 5. METHODS
s = add_slide(); set_bg(s, WHITE); title_bar(s, "Methods Implemented")
data = [
    ["Family", "Method", "Idea in one line"],
    ["Non-personalized", "Most Popular", "Rank by number of ratings"],
    ["Non-personalized", "Highest Average", "Top mean rating, min 20 ratings"],
    ["Non-personalized", "Random", "Lower-bound / coverage reference"],
    ["Content-based", "TF-IDF genres", "Cosine of user profile vs item genre vectors"],
    ["Collaborative", "Item-Item CF", "Adjusted cosine + shrinkage, k-NN over items"],
    ["Collaborative", "User-User CF", "Same, over the (denser) user axis"],
    ["Latent factor", "Matrix Factorization", "Biased SGD: mu + b_u + b_i + p_u·q_i"],
]
make_table(s, Inches(0.7), Inches(1.4), Inches(11.9), Inches(4.8), data, font=14)
tf = add_box(s, Inches(0.7), Inches(6.4), Inches(11.9), Inches(0.8))
bullets(tf, [
    ("MF is written from scratch in NumPy (scikit-surprise does not build on Python 3.14) — fully transparent.", 0, True, GREY),
], size=14, space=4)

# ============================================================ 6. PROTOCOL
s = add_slide(); set_bg(s, WHITE); title_bar(s, "Evaluation Protocol")
tf = add_box(s, Inches(0.6), Inches(1.3), Inches(12.1), Inches(5.6))
bullets(tf, [
    ("Split: per-user stratified 80/20 — every user keeps train history and held-out test items.", 0, True, ACCENT),
    ("Relevance: a test item counts as relevant if the user rated it ≥ 4.0 (positivity-aware).", 0, True, ACCENT),
    ("Top-N ranking metrics (K=10):", 0, True, ACCENT),
    ("Precision@K, Recall@K, NDCG@K, MRR@K, Hit-Rate@K.", 1),
    ("Beyond-accuracy metrics:", 0, True, ACCENT),
    ("Catalog coverage, novelty (mean −log₂ P(item)), intra-list diversity (1 − cosine over genres).", 1),
    ("Rating-prediction metrics (models with predict_score): RMSE, MAE.", 0, True, ACCENT),
    ("Averaged over the 599 users that have at least one relevant test item.", 0, False, GREY),
], size=19, space=10)

# ============================================================ 7. RESULTS TABLE
s = add_slide(); set_bg(s, WHITE); title_bar(s, "Method Comparison — Results (K=10)")
data = [
    ["Method", "NDCG", "P@10", "Recall", "HitRate", "Coverage", "Novelty", "Diversity"],
    ["User-User CF", "0.219", "0.163", "0.145", "0.648", "0.031", "2.19", "0.75"],
    ["Item-Item CF", "0.211", "0.155", "0.126", "0.644", "0.086", "2.80", "0.72"],
    ["Most Popular", "0.161", "0.122", "0.104", "0.574", "0.006", "1.66", "0.76"],
    ["Matrix Factorization", "0.072", "0.056", "0.041", "0.356", "0.049", "3.41", "0.78"],
    ["Highest Average", "0.062", "0.046", "0.034", "0.301", "0.004", "3.39", "0.66"],
    ["Content-Based", "0.008", "0.007", "0.005", "0.063", "0.341", "7.28", "0.085"],
    ["Random", "0.001", "0.001", "0.001", "0.013", "0.489", "7.57", "0.80"],
]
gt = make_table(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(4.4), data, font=13)
# highlight the best NDCG row
for c in range(len(data[0])):
    cell = gt.cell(1, c)
    for p in cell.text_frame.paragraphs:
        for run in p.runs:
            run.font.bold = True; run.font.color.rgb = WIN
tf = add_box(s, Inches(0.5), Inches(6.0), Inches(12.3), Inches(1.2))
bullets(tf, [
    ("Neighbourhood CF wins ranking; user-user edges item-item because the user axis is denser.", 0, True),
    ("Most-popular is a strong accuracy baseline but touches only 0.6% of the catalog — pure popularity bias.", 0),
], size=15, space=6)

# ============================================================ 8. ACCURACY VS COVERAGE
s = add_slide(); set_bg(s, WHITE); title_bar(s, "Accuracy is Not Enough")
add_image(s, "accuracy_vs_coverage.png", Inches(0.5), Inches(1.3), h=Inches(5.6))
tf = add_box(s, Inches(7.4), Inches(1.5), Inches(5.4), Inches(5.4))
bullets(tf, [
    ("The accuracy–coverage frontier", 0, True, ACCENT),
    ("Most accurate methods recommend from a tiny slice of the catalog.", 1),
    ("High-coverage methods (content, random) are inaccurate.", 1),
    ("No single winner", 0, True, ACCENT),
    ("A product must choose a point on this frontier.", 1),
    ("Discovery-focused service → accept lower accuracy for coverage/novelty.", 1),
    ("'Safe' service → popularity / CF.", 1),
], size=18, space=10)

# ============================================================ 9. NOVELTY / DIVERSITY
s = add_slide(); set_bg(s, WHITE); title_bar(s, "Beyond Accuracy: Novelty & Diversity")
add_image(s, "novelty_diversity.png", Inches(0.4), Inches(1.4), w=Inches(8.2))
tf = add_box(s, Inches(8.9), Inches(1.6), Inches(4.0), Inches(5.2))
bullets(tf, [
    ("Novelty: −log₂ P(item)", 0, True, ACCENT),
    ("Popular = low novelty; niche = high.", 1),
    ("Diversity: spread of a single list", 0, True, ACCENT),
    ("Content-based is highly novel but has very LOW diversity — it locks onto one genre cluster.", 1),
    ("These trade-offs are invisible to accuracy metrics.", 0, False, GREY),
], size=17, space=10)

# ============================================================ 10. MF NUANCE
s = add_slide(); set_bg(s, WHITE); title_bar(s, "Matrix Factorization: Right Tool, Wrong Metric")
add_image(s, "rating_prediction.png", Inches(0.5), Inches(1.4), w=Inches(6.6))
tf = add_box(s, Inches(7.4), Inches(1.4), Inches(5.4), Inches(5.6))
bullets(tf, [
    ("MF looks weak on top-N (NDCG 0.072)…", 0, True, ACCENT),
    ("…but it is the BEST rating predictor:", 0, True, WIN),
    ("RMSE 0.885 vs item-item 0.905; ~15% better than the global-mean baseline (1.04).", 1),
    ("Why the gap?", 0, True, ACCENT),
    ("MF minimizes squared error, not ranking. Good RMSE ≠ good top-10.", 1),
    ("Fix for production: ranking objective (BPR / implicit-feedback / logistic MF).", 1),
    ("'MF doesn't work' is wrong — it optimizes a different objective.", 0, False, GREY),
], size=17, space=9)

# ============================================================ 10b. ABLATION TF-IDF
s = add_slide(); set_bg(s, WHITE); title_bar(s, "Ablation — TF-IDF vs. Raw Genre Vectors")
tf = add_box(s, Inches(0.6), Inches(1.2), Inches(12.1), Inches(0.8))
bullets(tf, [
    ("Does weighting genres by rarity (TF-IDF) beat treating every genre equally (raw 0/1)?", 0, True),
], size=16)
# Same metric set as the main results slide; winner of each column highlighted.
data = [
    ["Content-based variant", "NDCG", "P@10", "Recall", "HitRate", "Coverage", "Novelty", "Diversity"],
    ["TF-IDF genres", "0.0076", "0.0068", "0.0047", "0.063", "0.341", "7.28", "0.085"],
    ["Raw genre vectors", "0.0056", "0.0058", "0.0035", "0.053", "0.352", "7.29", "0.075"],
]
gt = make_table(s, Inches(0.5), Inches(2.1), Inches(12.3), Inches(1.5), data, font=13)
# Higher is better for every column here; bold+green the winning cell per metric.
winner_row = [1, 1, 1, 1, 2, 2, 1]  # 1 = TF-IDF, 2 = Raw, for cols 1..7
for c, wr in enumerate(winner_row, start=1):
    cell = gt.cell(wr, c)
    for p in cell.text_frame.paragraphs:
        for run in p.runs:
            run.font.bold = True; run.font.color.rgb = WIN
tf = add_box(s, Inches(0.6), Inches(4.1), Inches(12.1), Inches(2.8))
bullets(tf, [
    ("TF-IDF wins 5 of 7 metrics — all four accuracy metrics plus diversity.", 0, True, ACCENT),
    ("NDCG +36% relative (0.0056 → 0.0076): IDF up-weights distinctive genres (Film-Noir, Western), building more discriminative profiles.", 1),
    ("Raw genre vectors only edge ahead on coverage and novelty.", 0, True, ACCENT),
    ("Treating every genre equally spreads picks slightly wider across the catalog.", 1),
    ("The effect is small by construction:", 0, True, ACCENT),
    ("~20 genre tokens and term-frequency ≡ 1, so only the IDF term is active — richer text (tags, synopsis) would help more.", 1),
], size=17, space=9)

# ============================================================ 11. REC EXAMPLES
s = add_slide(); set_bg(s, WHITE); title_bar(s, "Recommendation Examples (User 599)")
tf = add_box(s, Inches(0.6), Inches(1.2), Inches(12.1), Inches(0.7))
bullets(tf, [("Same user, top-5 from each method — the personalities are clearly different:", 0, True)], size=16)
data = [
    ["Method", "Top recommendations"],
    ["Most Popular", "Matrix · Star Wars IV · Schindler's List · Empire Strikes Back"],
    ["Item-Item CF", "Eternal Sunshine · Matrix · Kill Bill 2 · Star Wars IV"],
    ["User-User CF", "Star Wars IV · Matrix · Empire Strikes Back · LotR Fellowship"],
    ["Content-Based", "Home for the Holidays · Mr. Holland's Opus · Tape · Focus"],
    ["Matrix Fact.", "Lawrence of Arabia · Donnie Darko · Kill Bill 2 · Some Like It Hot"],
]
make_table(s, Inches(0.6), Inches(2.0), Inches(12.1), Inches(3.6), data, font=14)
tf = add_box(s, Inches(0.6), Inches(5.9), Inches(12.1), Inches(1.0))
bullets(tf, [
    ("CF leans to well-loved classics; content-based surfaces niche same-genre titles; MF mixes acclaimed films.", 0, False, GREY),
], size=14, space=4)

# ============================================================ 12. CHALLENGES
s = add_slide(); set_bg(s, WHITE); title_bar(s, "Technical Challenges")
tf = add_box(s, Inches(0.6), Inches(1.3), Inches(12.1), Inches(5.6))
bullets(tf, [
    ("Ranking vs. rating in CF", 0, True, ACCENT),
    ("Normalizing scores by Σ|sim| let single-neighbour tail items flood top-N (P@10 collapsed to ~0.006).", 1),
    ("Fix: rank by the raw weighted sum; keep normalization only for rating prediction. P@10 → 0.155.", 1),
    ("Tuning MF for top-N", 0, True, ACCENT),
    ("More regularization helped RMSE but hurt ranking; settled on 50 factors, 30 epochs, λ=0.01.", 1),
    ("Library availability", 0, True, ACCENT),
    ("scikit-surprise won't build on Python 3.14 → implemented biased SGD MF in pure NumPy.", 1),
    ("Engineering", 0, True, ACCENT),
    ("Kept matrices sparse; pre-trained & pickled models so the Streamlit app loads instantly; UTF-8 console fix.", 1),
], size=17, space=7)

# ============================================================ 13. PROTOTYPE / UI
s = add_slide(); set_bg(s, WHITE); title_bar(s, "The Prototype (Streamlit)")
tf = add_box(s, Inches(0.6), Inches(1.3), Inches(12.1), Inches(5.6))
bullets(tf, [
    ("Explore tab", 0, True, ACCENT),
    ("Pick a user → see their taste profile (top genres, favourite films).", 1),
    ("Compare any subset of algorithms side by side for that user.", 1),
    ("Every recommendation carries a 'why': matching genres, a similar film you liked, predicted rating, or popularity.", 1),
    ("Evaluation tab", 0, True, ACCENT),
    ("The full metrics table and comparison figures, in-app.", 1),
    ("Engineering", 0, True, ACCENT),
    ("Models pre-trained via build_artifacts.py and cached → instant, responsive UX.", 1),
    ("Run:  python build_artifacts.py    →    streamlit run app.py", 0, True, GREY),
], size=18, space=8)

# ============================================================ 14. FINAL REMARKS
s = add_slide(); set_bg(s, WHITE); title_bar(s, "Final Remarks & Future Work")
tf = add_box(s, Inches(0.6), Inches(1.3), Inches(12.1), Inches(5.6))
bullets(tf, [
    ("What the comparison taught us", 0, True, ACCENT),
    ("Neighbourhood CF gives the best top-N here; popularity is a deceptively strong, low-coverage baseline.", 1),
    ("Accuracy and beyond-accuracy goals genuinely conflict — the right model depends on the product.", 1),
    ("Matrix factorization wins rating prediction; ranking needs a ranking objective.", 1),
    ("Future work", 0, True, ACCENT),
    ("Hybrid CF + content to cover the cold/long-tail items CF cannot score.", 1),
    ("Implicit/ranking-loss MF (BPR, ALS) to make latent factors competitive on top-N.", 1),
    ("Temporal split, cross-validation, and explicit cold-start analysis.", 1),
    ("Thank you.", 0, True, WIN),
], size=18, space=8)

prs.save(OUT)
print(f"Saved deck with {len(prs.slides)} slides to {OUT}")
